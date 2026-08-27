from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt
from lxml import etree
from copy import deepcopy

INK = RGBColor(0x2F, 0x2A, 0x26)
BONE = RGBColor(0xF7, 0xF4, 0xEF)
WHITE = RGBColor(0xFF, 0xFD, 0xF9)
ACCENT = RGBColor(0xD4, 0x78, 0x5A)
ACCENT_DARK = RGBColor(0xB8, 0x5C, 0x42)
UMBER = RGBColor(0x6F, 0x62, 0x59)
PARCHMENT = RGBColor(0xED, 0xE7, 0xDC)
KHAKI = RGBColor(0xC9, 0xB8, 0xA8)

W, H = Inches(13.333), Inches(7.5)


def set_run(run, text, size, bold=False, color=INK, name="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    run.font.italic = False


def add_text(tf, text, size, bold=False, color=INK, align=PP_ALIGN.LEFT, space_after=6):
    p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    set_run(run, text, size, bold, color)
    return p


def rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def round_rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.adjustments[0] = 0.08
    return s


def tb(slide, l, t, w, h):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    return tf


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, BONE)
    rect(slide, 0, 0, Inches(0.18), H, ACCENT)
    return slide


def kicker(slide, text):
    tf = tb(slide, Inches(0.7), Inches(0.38), Inches(10), Inches(0.4))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_run(run, text.upper(), 12, True, ACCENT_DARK)
    run.font.name = "Calibri"


def title(slide, text, top=0.75):
    tf = tb(slide, Inches(0.7), Inches(top), Inches(12), Inches(1.1))
    p = tf.paragraphs[0]
    run = p.add_run()
    set_run(run, text, 32, True, INK)


def bullets(slide, items, top=2.05, left=0.7, width=12.0, size=20):
    tf = tb(slide, Inches(left), Inches(top), Inches(width), Inches(4.8))
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = 0
        p.space_after = Pt(14)
        p.line_spacing = 1.15
        run = p.add_run()
        set_run(run, item, size, False, INK)


def card(slide, l, t, w, h, heading, body):
    round_rect(slide, Inches(l), Inches(t), Inches(w), Inches(h), WHITE)
    tf = tb(slide, Inches(l + 0.28), Inches(t + 0.28), Inches(w - 0.5), Inches(h - 0.5))
    p = tf.paragraphs[0]
    run = p.add_run()
    set_run(run, heading, 18, True, ACCENT_DARK)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(10)
    p2.line_spacing = 1.15
    run2 = p2.add_run()
    set_run(run2, body, 15, False, UMBER)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1 Portada
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, INK)
    rect(s, 0, 0, Inches(0.22), H, ACCENT)
    tf = tb(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(0.4))
    p = tf.paragraphs[0]
    run = p.add_run()
    set_run(run, "APLICACIONES MÓVILES  ·  UPB  ·  ENTREGA 2", 14, True, KHAKI)
    tf = tb(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(1.4))
    p = tf.paragraphs[0]
    run = p.add_run()
    set_run(run, "UPVenta", 60, True, WHITE)
    tf = tb(s, Inches(0.9), Inches(3.35), Inches(11.2), Inches(1.1))
    p = tf.paragraphs[0]
    p.line_spacing = 1.2
    run = p.add_run()
    set_run(run, "Inventario y ventas en el celular,\npara vendedores de la Universidad Pontificia Bolivariana.", 24, False, BONE)
    tf = tb(s, Inches(0.9), Inches(5.35), Inches(11), Inches(1.3))
    p = tf.paragraphs[0]
    run = p.add_run()
    set_run(run, "Autores", 13, True, KHAKI)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    run = p2.add_run()
    set_run(run, "Nicolás Agudelo", 22, True, WHITE)
    p3 = tf.add_paragraph()
    p3.space_before = Pt(4)
    run = p3.add_run()
    set_run(run, "Sebastian Muñoz", 22, True, WHITE)

    # 2 Problema
    s = blank(prs)
    kicker(s, "El problema")
    title(s, "Vender en la UPB, sin perder el control")
    bullets(
        s,
        [
            "En el campus hay decenas de puestos, cafeterías y emprendimientos: el volumen de ventas es alto y rápido.",
            "Muchos anotan en papel, WhatsApp o de memoria. Al cierre del día no saben cuánto vendieron ni qué se acabó.",
            "Si hay más de una persona atendiendo, no hay una cuenta clara: cualquiera mezcla el inventario y se pierde el recaudo.",
            "El resultado: quiebres de stock, plata que no cuadra y tiempo perdido reconstruyendo la jornada.",
        ],
        top=2.0,
        size=20,
    )

    # 3 Solución
    s = blank(prs)
    kicker(s, "La solución")
    title(s, "UPVenta: la caja en el bolsillo")
    tf = tb(s, Inches(0.7), Inches(1.95), Inches(12), Inches(0.8))
    p = tf.paragraphs[0]
    p.line_spacing = 1.2
    run = p.add_run()
    set_run(
        run,
        "Una app híbrida para que el vendedor registre la venta en el celular apenas ocurre, y el inventario se actualice solo.",
        20,
        False,
        UMBER,
    )
    card(
        s, 0.7, 3.05, 3.8, 3.4,
        "A escala campus",
        "Pensada para muchos puestos a la vez: cada negocio entra con su cuenta, su equipo y su catálogo. No depende de un computador ni de internet.",
    )
    card(
        s, 4.75, 3.05, 3.8, 3.4,
        "En el momento de la venta",
        "Se toca +, se confirma, y queda el ticket: producto, cantidad, precio y stock que quedó. La relación venta ↔️ producto no se pierde.",
    )
    card(
        s, 8.8, 3.05, 3.8, 3.4,
        "Cierre del día",
        "Recaudo, ticket promedio, más vendidos y alertas de stock bajo. El vendedor cierra caja sin armar una hoja de cálculo.",
    )

    # 4 User persona
    s = blank(prs)
    kicker(s, "User persona")
    title(s, "Camila, estudiante que busca ingresos extra")
    round_rect(s, Inches(0.7), Inches(2.0), Inches(4.35), Inches(4.55), WHITE)
    tf = tb(s, Inches(0.98), Inches(2.22), Inches(3.85), Inches(4.15))
    p = tf.paragraphs[0]
    run = p.add_run()
    set_run(run, "Camila Restrepo", 22, True, INK)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    run = p2.add_run()
    set_run(run, "21 años  ·  Estudiante UPB", 15, False, ACCENT_DARK)
    p3 = tf.add_paragraph()
    p3.space_before = Pt(4)
    run = p3.add_run()
    set_run(run, "Dulce & Tela  ·  brownies y accesorios", 14, False, UMBER)
    p4 = tf.add_paragraph()
    p4.space_before = Pt(16)
    p4.line_spacing = 1.2
    run = p4.add_run()
    set_run(
        run,
        "“Entre clase y clase armo un puesto. Quiero saber cuánto gané hoy y qué se me acabó, sin llevar un cuaderno.”",
        16,
        False,
        INK,
    )
    card(
        s, 5.3, 2.0, 7.2, 1.35,
        "Meta",
        "Sacar un ingreso extra en el campus sin descuidar el semestre: vender rápido, cuadrar caja y no quedarse sin producto.",
    )
    card(
        s, 5.3, 3.5, 7.2, 1.4,
        "Frustración",
        "Anota ventas en WhatsApp o de memoria. En el pico del descanso se le olvida restar stock y al final no sabe si le alcanzó para reponer.",
    )
    card(
        s, 5.3, 5.05, 7.2, 1.5,
        "Qué necesita de UPVenta",
        "Entrar con su usuario, cargar el catálogo una vez y registrar cada venta en el celular. Alertas de stock bajo y un resumen del día.",
    )

    # 5 Productos
    s = blank(prs)
    kicker(s, "Gestión de productos")
    title(s, "El catálogo vive en el celular")
    bullets(
        s,
        [
            "Alta y edición de productos: nombre, categoría, precio, unidad, stock y umbral mínimo.",
            "Búsqueda y filtros (postres, ropa, accesorios…) para encontrar rápido en un catálogo grande.",
            "Cuando el stock llega al mínimo, la campana avisa y se puede reponer 1 o 5 unidades al toque.",
            "Cada venta descuenta inventario: no hay que restar a mano ni al final del día.",
            "Todo queda ligado por id: el historial muestra qué se vendió y cuánto stock quedó después.",
        ],
        top=1.95,
        size=19,
    )

    # 6 Flujo
    s = blank(prs)
    kicker(s, "Cómo se usa en el puesto")
    title(s, "De la venta al cierre, en cuatro toques")
    steps = [
        ("1", "Vender", "Se suma la cantidad apenas el cliente paga. El stock se descuenta al confirmar."),
        ("2", "Hoy", "Recaudo, número de ventas, unidades y ticket promedio del día."),
        ("3", "Alertas", "Lo que está por acabarse, para no decir “ya no hay” a media jornada."),
        ("4", "Historial", "Cada ticket con detalle. Se puede volver a una venta y ver productos y cantidades."),
    ]
    for i, (num, head, body) in enumerate(steps):
        x = 0.7 + i * 3.1
        round_rect(s, Inches(x), Inches(2.15), Inches(2.9), Inches(4.2), WHITE)
        nbox = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.28), Inches(2.4), Inches(0.48), Inches(0.48))
        nbox.fill.solid()
        nbox.fill.fore_color.rgb = ACCENT
        nbox.line.fill.background()
        ntf = nbox.text_frame
        ntf.word_wrap = False
        ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = ntf.paragraphs[0].add_run()
        set_run(run, num, 16, True, WHITE)
        tf = tb(s, Inches(x + 0.28), Inches(3.1), Inches(2.35), Inches(2.9))
        p = tf.paragraphs[0]
        run = p.add_run()
        set_run(run, head, 20, True, INK)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(10)
        p2.line_spacing = 1.2
        run = p2.add_run()
        set_run(run, body, 15, False, UMBER)

    # 7 Cierre
    s = blank(prs)
    kicker(s, "Para la demo")
    title(s, "Qué queremos que se lleven")
    bullets(
        s,
        [
            "UPVenta le quita el cuaderno al vendedor de la UPB: equipo con login, productos con stock y venta al instante.",
            "Escala a muchos puestos porque cada negocio es independiente y funciona sin internet (localStorage).",
            "Stack: SPA Vanilla JS + Vite + SASS. Un solo index.html, lista para celular.",
            "Prototipo en Figma + app. Demo: demo@upventa.app  /  demo123. Preguntas.",
        ],
        top=1.95,
        size=20,
    )
    tf = tb(s, Inches(0.7), Inches(6.35), Inches(12), Inches(0.6))
    p = tf.paragraphs[0]
    run = p.add_run()
    set_run(run, "Nicolás Agudelo  ·  Sebastian Muñoz  ·  Aplicaciones móviles — UPB", 16, False, UMBER)

    out = "/Users/niscko/Desktop/Apps moviles/presentacion/UPVenta-exposicion.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    print(build())
